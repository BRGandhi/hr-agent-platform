from __future__ import annotations

import re
from dataclasses import dataclass
from datetime import datetime
from io import BytesIO

import pandas as pd
import xlsxwriter
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ACCENT = (11, 92, 171)
ACCENT_DARK = (8, 52, 102)
ACCENT_SOFT = (236, 243, 250)
INK = (24, 39, 75)
MUTED = (92, 108, 132)
LINE = (209, 219, 231)
PANEL = (247, 250, 253)
WHITE = (255, 255, 255)
SUCCESS = (15, 118, 110)
ALERT = (192, 60, 60)
CHART_COLORS = [
    (11, 92, 171),
    (46, 125, 202),
    (92, 176, 217),
    (15, 118, 110),
    (230, 126, 34),
    (192, 57, 43),
]
MEASURE_KEYWORDS = (
    "count",
    "total",
    "rate",
    "pct",
    "percent",
    "avg",
    "average",
    "mean",
    "score",
    "salary",
    "income",
    "pay",
    "tenure",
    "years",
    "employees",
    "headcount",
    "attrited",
    "promotion",
)


@dataclass
class ReportMetric:
    label: str
    value: str
    note: str = ""


@dataclass
class ReportStory:
    title: str
    role: str
    scope_name: str
    generated_label: str
    headline: str
    subheadline: str
    key_metrics: list[ReportMetric]
    insights: list[str]
    actions: list[str]
    chart_title: str
    chart_subtitle: str
    chart_spec: dict
    chart_image: bytes | None
    preview_rows: list[dict]
    source_note: str


def sanitize_filename(title: str, extension: str) -> str:
    cleaned = re.sub(r"[^a-z0-9]+", "_", str(title or "report").strip().lower()).strip("_")
    return f"{cleaned or 'report'}.{extension.lstrip('.')}"


def configure_export_rows(
    rows: list[dict],
    *,
    columns: list[str] | None = None,
    sort_by: str = "",
    sort_direction: str = "asc",
    max_rows: int | None = None,
    filter_column: str = "",
    filter_value: str = "",
) -> list[dict]:
    df = _prepare_dataframe(rows)
    if df.empty:
        return []

    if filter_column and filter_column in df.columns and str(filter_value or "").strip():
        df = df[df[filter_column].astype(str) == str(filter_value)]

    if sort_by and sort_by in df.columns:
        ascending = str(sort_direction or "asc").lower() != "desc"
        numeric_sort = pd.to_numeric(df[sort_by], errors="coerce")
        if numeric_sort.notna().sum() >= max(3, int(len(df) * 0.6)):
            df = df.assign(__sort_value=numeric_sort).sort_values("__sort_value", ascending=ascending, na_position="last")
            df = df.drop(columns="__sort_value")
        else:
            df = df.sort_values(sort_by, ascending=ascending, na_position="last")

    if columns:
        valid_columns = [column for column in columns if column in df.columns]
        if valid_columns:
            df = df[valid_columns]

    if max_rows and int(max_rows) > 0:
        df = df.head(int(max_rows))

    return df.fillna("").to_dict(orient="records")


def build_configured_excel(
    title: str,
    rows: list[dict],
    *,
    columns: list[str] | None = None,
    sort_by: str = "",
    sort_direction: str = "asc",
    max_rows: int | None = None,
    include_summary: bool = True,
    filter_column: str = "",
    filter_value: str = "",
) -> bytes:
    configured_rows = configure_export_rows(
        rows,
        columns=columns,
        sort_by=sort_by,
        sort_direction=sort_direction,
        max_rows=max_rows,
        filter_column=filter_column,
        filter_value=filter_value,
    )
    selected_columns = (
        list(configured_rows[0].keys()) if configured_rows
        else list(rows[0].keys()) if rows
        else ["Message"]
    )

    buffer = BytesIO()
    workbook = xlsxwriter.Workbook(buffer, {"in_memory": True})
    header = workbook.add_format({"bold": True, "bg_color": "#EEF5FB", "font_color": "#0B3A75", "border": 1})
    cell = workbook.add_format({"font_color": "#334155", "border": 1})
    note = workbook.add_format({"font_color": "#64748B"})

    if include_summary:
        summary = workbook.add_worksheet("Summary")
        summary.set_column("A:A", 22)
        summary.set_column("B:B", 46)
        summary.write("A1", "Report", header)
        summary.write("B1", title or "HR export", cell)
        summary.write("A2", "Rows exported", header)
        summary.write("B2", len(configured_rows), cell)
        summary.write("A3", "Columns", header)
        summary.write("B3", ", ".join(selected_columns), cell)
        summary.write("A4", "Filter", header)
        summary.write("B4", f"{filter_column} = {filter_value}" if filter_column and filter_value else "None", cell)
        summary.write("A5", "Sort", header)
        summary.write("B5", f"{sort_by or 'None'} ({sort_direction})", cell)
        summary.write("A6", "Generated", header)
        summary.write("B6", datetime.now().strftime("%B %d, %Y %I:%M %p"), cell)
        summary.write("A8", "Note", header)
        summary.write("B8", "This workbook reflects the current governed view and the builder selections applied at export time.", note)

    worksheet = workbook.add_worksheet(_sheet_name(title or "Export"))
    worksheet.freeze_panes(1, 0)
    worksheet.autofilter(0, 0, max(len(configured_rows), 1), max(len(selected_columns) - 1, 0))

    for column_index, column in enumerate(selected_columns):
        worksheet.write(0, column_index, column, header)
        worksheet.set_column(column_index, column_index, max(14, min(28, len(str(column)) + 4)))

    if configured_rows:
        for row_index, row in enumerate(configured_rows, start=1):
            for column_index, column in enumerate(selected_columns):
                worksheet.write(row_index, column_index, _excel_safe_value(row.get(column)), cell)
    else:
        worksheet.write(1, 0, "No rows matched the current selections.", note)

    workbook.close()
    buffer.seek(0)
    return buffer.getvalue()


def build_report_story(
    title: str,
    rows: list[dict],
    *,
    role: str,
    scope_name: str,
    question: str = "",
    chart_spec: dict | None = None,
) -> ReportStory:
    df = _prepare_dataframe(rows)
    chart_spec = chart_spec or {}
    metric_col = _choose_metric_column(df)
    dimension_col = _choose_dimension_column(df, exclude={metric_col} if metric_col else None)
    chart_spec = {
        "chart_type": str(chart_spec.get("chart_type", "") or _default_chart_type(df, dimension_col, metric_col)),
        "x_column": str(chart_spec.get("x_column", "") or dimension_col),
        "y_column": str(chart_spec.get("y_column", "") or metric_col),
        "color_column": str(chart_spec.get("color_column", "") or ""),
        "title": str(chart_spec.get("title", "") or _default_chart_title(dimension_col, metric_col)),
    }
    return ReportStory(
        title=title or "HR insight brief",
        role=role or "HR partner",
        scope_name=scope_name or "Approved scope",
        generated_label=datetime.now().strftime("%B %d, %Y"),
        headline=_build_headline(df, dimension_col, metric_col, scope_name),
        subheadline=_build_subheadline(df, dimension_col, metric_col, question),
        key_metrics=_build_key_metrics(df, dimension_col, metric_col),
        insights=_build_insights(df, dimension_col, metric_col),
        actions=_build_actions(dimension_col),
        chart_title=chart_spec["title"],
        chart_subtitle=_chart_subtitle(chart_spec, dimension_col, metric_col),
        chart_spec=chart_spec,
        chart_image=None,
        preview_rows=list(rows[:5]),
        source_note=_source_note(rows, scope_name),
    )


def attach_story_chart(story: ReportStory, rows: list[dict]) -> ReportStory:
    story.chart_image = render_chart_image(rows, story.chart_spec)
    return story


def _prepare_dataframe(rows: list[dict]) -> pd.DataFrame:
    df = pd.DataFrame(rows or []).copy()
    if df.empty:
        return df
    for column in df.columns:
        series = df[column]
        if pd.api.types.is_numeric_dtype(series) or pd.api.types.is_datetime64_any_dtype(series):
            continue
        numeric = pd.to_numeric(series, errors="coerce")
        if numeric.notna().sum() >= max(3, int(series.notna().sum() * 0.75)):
            df[column] = numeric
            continue
        if _looks_temporal(column):
            parsed = pd.to_datetime(series, errors="coerce")
            if parsed.notna().sum() >= max(2, int(series.notna().sum() * 0.65)):
                df[column] = parsed
    return df


def _sheet_name(title: str) -> str:
    cleaned = "".join("_" if char in '[]:*?/\\' else char for char in str(title or "Export")).strip()
    return (cleaned or "Export")[:31]


def _source_note(rows: list[dict], scope_name: str) -> str:
    if rows and "SnapshotMonth" in (rows[0] or {}):
        count = f"{len(rows):,} governed monthly rows" if rows else "governed monthly rows"
        return (
            f"Source: {count} from the simulated workforce trend layer for the current {scope_name} HR view."
        )
    if rows:
        return f"Source: {len(rows):,} governed rows from the current {scope_name} HR view."
    return f"Source: governed {scope_name} HR view."


def _excel_safe_value(value):
    if pd.isna(value):
        return ""
    if isinstance(value, pd.Timestamp):
        return value.to_pydatetime().replace(microsecond=0).isoformat()
    return value


def _choose_metric_column(df: pd.DataFrame, exclude: set[str] | None = None) -> str:
    exclude = exclude or set()
    numeric_columns = [column for column in df.columns if column not in exclude and pd.api.types.is_numeric_dtype(df[column])]
    if not numeric_columns:
        return ""
    scored: list[tuple[float, str]] = []
    for column in numeric_columns:
        score = min(df[column].nunique(dropna=True), 8) / 2
        if any(token in column.lower() for token in MEASURE_KEYWORDS):
            score += 4
        if df[column].nunique(dropna=True) <= 1:
            score -= 8
        scored.append((score, column))
    return max(scored)[1]


def _choose_dimension_column(df: pd.DataFrame, exclude: set[str] | None = None) -> str:
    exclude = exclude or set()
    scored: list[tuple[float, str]] = []
    for column in df.columns:
        if column in exclude:
            continue
        series = df[column]
        if pd.api.types.is_numeric_dtype(series):
            continue
        unique_values = series.dropna().nunique()
        if unique_values <= 1:
            continue
        score = 0.0
        if 2 <= unique_values <= 8:
            score += 5
        elif unique_values <= 14:
            score += 3
        elif unique_values <= 20:
            score += 1
        if any(token in column.lower() for token in ("department", "job", "role", "level", "gender", "travel", "status", "balance")):
            score += 2
        scored.append((score, column))
    return max(scored)[1] if scored else ""


def _default_chart_type(df: pd.DataFrame, dimension_col: str, metric_col: str) -> str:
    if dimension_col and metric_col:
        if pd.api.types.is_datetime64_any_dtype(df[dimension_col]):
            return "line"
        unique_values = df[dimension_col].dropna().nunique()
        if 2 <= unique_values <= 6 and any(token in str(metric_col).lower() for token in ("pct", "rate", "percent", "share")):
            return "lollipop"
        return "lollipop"
    if metric_col:
        return "indicator"
    return "summary"


def _default_chart_title(dimension_col: str, metric_col: str) -> str:
    if dimension_col and metric_col:
        return f"{_humanize(metric_col)} by {_humanize(dimension_col)}"
    if metric_col:
        return _humanize(metric_col)
    return "HR summary"


def _build_headline(df: pd.DataFrame, dimension_col: str, metric_col: str, scope_name: str) -> str:
    if df.empty:
        return f"{scope_name} insight brief is ready."
    temporal_frame = _time_series_frame(df, dimension_col, metric_col)
    if not temporal_frame.empty:
        latest = temporal_frame.iloc[-1]
        latest_label = _format_period_point(latest[dimension_col])
        if "YoYHeadcountChangePct" in df.columns and pd.notna(latest.get("YoYHeadcountChangePct")):
            return (
                f"{_humanize(metric_col)} ended at {_format_value(metric_col, latest[metric_col])} in {latest_label}, "
                f"with {_format_delta_value('YoYHeadcountChangePct', latest.get('YoYHeadcountChangePct'))} versus last year."
            )
        if len(temporal_frame) >= 2:
            prior = temporal_frame.iloc[-2]
            delta = float(latest[metric_col]) - float(prior[metric_col])
            return (
                f"{_humanize(metric_col)} ended at {_format_value(metric_col, latest[metric_col])} in {latest_label}, "
                f"moving {_format_delta_value(metric_col, delta)} versus the prior month."
            )
    if dimension_col and metric_col:
        frame = df[[dimension_col, metric_col]].dropna().sort_values(metric_col, ascending=False)
        if not frame.empty:
            top_row = frame.iloc[0]
            return f"{top_row[dimension_col]} is leading on {_humanize(metric_col).lower()} in the current view."
    if metric_col:
        return f"{_humanize(metric_col)} is the clearest signal in the current {scope_name} slice."
    return f"{scope_name} insight brief is ready for review."


def _build_subheadline(df: pd.DataFrame, dimension_col: str, metric_col: str, question: str) -> str:
    if df.empty:
        return "No governed rows are available in the current view."
    temporal_frame = _time_series_frame(df, dimension_col, metric_col)
    if not temporal_frame.empty:
        latest = temporal_frame.iloc[-1]
        start = temporal_frame.iloc[0]
        start_label = _format_period_point(start[dimension_col])
        latest_label = _format_period_point(latest[dimension_col])
        return (
            f"The period runs from {start_label} to {latest_label}. "
            f"{_humanize(metric_col)} moved from {_format_value(metric_col, start[metric_col])} "
            f"to {_format_value(metric_col, latest[metric_col])} across the selected reporting window."
        )
    if dimension_col and metric_col:
        frame = df[[dimension_col, metric_col]].dropna().sort_values(metric_col, ascending=False)
        if not frame.empty:
            top_row = frame.iloc[0]
            bottom_row = frame.iloc[-1]
            return (
                f"{_humanize(metric_col)} ranges from {_format_value(metric_col, bottom_row[metric_col])} "
                f"to {_format_value(metric_col, top_row[metric_col])} across the visible {_humanize(dimension_col).lower()} groups."
            )
    if metric_col:
        series = pd.to_numeric(df[metric_col], errors="coerce").dropna()
        if not series.empty:
            return (
                f"The current view spans {len(df):,} governed rows with an average {_humanize(metric_col).lower()} "
                f"of {_format_value(metric_col, series.mean())}."
            )
    return question or f"The current view includes {len(df):,} governed rows."


def _build_key_metrics(df: pd.DataFrame, dimension_col: str, metric_col: str) -> list[ReportMetric]:
    temporal_frame = _time_series_frame(df, dimension_col, metric_col)
    if not temporal_frame.empty:
        latest = temporal_frame.iloc[-1]
        start = temporal_frame.iloc[0]
        metrics = [
            ReportMetric(
                f"Latest {_humanize(metric_col)}",
                _format_value(metric_col, latest[metric_col]),
                _format_period_point(latest[dimension_col]),
            ),
            ReportMetric(
                "Period change",
                _format_delta_value(metric_col, float(latest[metric_col]) - float(start[metric_col])),
                f"{_format_period_point(start[dimension_col])} to {_format_period_point(latest[dimension_col])}",
            ),
        ]
        if "MoMHeadcountChangePct" in df.columns and pd.notna(latest.get("MoMHeadcountChangePct")):
            metrics.append(ReportMetric("MoM change", _format_delta_value("MoMHeadcountChangePct", latest.get("MoMHeadcountChangePct")), "Latest month"))
        elif "YoYHeadcountChangePct" in df.columns and pd.notna(latest.get("YoYHeadcountChangePct")):
            metrics.append(ReportMetric("YoY change", _format_delta_value("YoYHeadcountChangePct", latest.get("YoYHeadcountChangePct")), "Latest month"))
        while len(metrics) < 3:
            metrics.append(ReportMetric("Scope", "Governed", "Role-based access"))
        return metrics[:3]

    metrics = [ReportMetric("Rows in view", f"{len(df):,}", "Governed result set")]
    if metric_col and metric_col in df.columns:
        series = pd.to_numeric(df[metric_col], errors="coerce").dropna()
        if not series.empty:
            metrics.append(ReportMetric(_humanize(metric_col), _format_value(metric_col, series.sum()), "Aggregate total"))
            metrics.append(ReportMetric(f"Average {_humanize(metric_col).lower()}", _format_value(metric_col, series.mean()), "Across the current rows"))
    if len(metrics) < 3 and dimension_col and metric_col:
        frame = df[[dimension_col, metric_col]].dropna().sort_values(metric_col, ascending=False)
        if not frame.empty:
            metrics.append(ReportMetric("Top group", str(frame.iloc[0][dimension_col]), _format_value(metric_col, frame.iloc[0][metric_col])))
    while len(metrics) < 3:
        metrics.append(ReportMetric("Scope", "Governed", "Role-based access"))
    return metrics[:3]


def _build_insights(df: pd.DataFrame, dimension_col: str, metric_col: str) -> list[str]:
    if df.empty:
        return ["The current view is empty, so the export is acting as a placeholder rather than an analytical brief."]
    temporal_frame = _time_series_frame(df, dimension_col, metric_col)
    if not temporal_frame.empty:
        latest = temporal_frame.iloc[-1]
        start = temporal_frame.iloc[0]
        peak = temporal_frame.loc[temporal_frame[metric_col].idxmax()]
        trough = temporal_frame.loc[temporal_frame[metric_col].idxmin()]
        insights = [
            (
                f"{_humanize(metric_col)} moved from {_format_value(metric_col, start[metric_col])} "
                f"in {_format_period_point(start[dimension_col])} to {_format_value(metric_col, latest[metric_col])} "
                f"in {_format_period_point(latest[dimension_col])}."
            ),
            (
                f"The peak month was {_format_period_point(peak[dimension_col])} at "
                f"{_format_value(metric_col, peak[metric_col])}, while the low point was "
                f"{_format_period_point(trough[dimension_col])} at {_format_value(metric_col, trough[metric_col])}."
            ),
        ]
        if "MoMHeadcountChangePct" in df.columns and pd.notna(latest.get("MoMHeadcountChangePct")):
            insights.append(
                f"The latest month closed at {_format_delta_value('MoMHeadcountChangePct', latest.get('MoMHeadcountChangePct'))} versus the prior month."
            )
        elif "YoYHeadcountChangePct" in df.columns and pd.notna(latest.get("YoYHeadcountChangePct")):
            insights.append(
                f"The latest month closed at {_format_delta_value('YoYHeadcountChangePct', latest.get('YoYHeadcountChangePct'))} versus the prior year."
            )
        else:
            insights.append("This is a time-series view, so direction and turning points matter more than a single snapshot.")
        return insights
    if dimension_col and metric_col:
        frame = df[[dimension_col, metric_col]].dropna().sort_values(metric_col, ascending=False)
        if not frame.empty:
            top_row = frame.iloc[0]
            bottom_row = frame.iloc[-1]
            total = pd.to_numeric(frame[metric_col], errors="coerce").sum()
            share = (float(top_row[metric_col]) / total * 100.0) if total else 0.0
            return [
                f"{top_row[dimension_col]} is highest at {_format_value(metric_col, top_row[metric_col])}, compared with {_format_value(metric_col, bottom_row[metric_col])} for {bottom_row[dimension_col]}.",
                f"The leading group accounts for {share:.1f}% of the measured total in the current cut." if total else f"{len(frame):,} groups are visible in the current cut.",
                "This is a prioritization view, so the export is optimized for what should be reviewed first.",
            ]
    if metric_col:
        series = pd.to_numeric(df[metric_col], errors="coerce").dropna()
        if not series.empty:
            return [
                f"The measure spans from {_format_value(metric_col, series.min())} to {_format_value(metric_col, series.max())}.",
                f"The median sits at {_format_value(metric_col, series.median())}, which helps anchor what normal looks like.",
                "A distribution or indicator view is stronger than a simple category chart for this result set.",
            ]
    return [
        "The current result set is more descriptive than metric-heavy.",
        "The main value here is guiding the next HR question rather than over-explaining the current cut.",
        "Use the configured Excel export if the audience needs row-level detail.",
    ]


def _build_actions(dimension_col: str) -> list[str]:
    if dimension_col and _looks_temporal(dimension_col):
        return [
            "Review the inflection points first, not just the latest month.",
            "Pressure-test which departments or roles are driving the latest movement.",
            "Use the period trend as the lead story and the roster table as backup detail.",
        ]
    if dimension_col:
        return [
            f"Review the leading {_humanize(dimension_col).lower()} group first.",
            f"Pressure-test whether the pattern holds when you cut the data by department, role, and tenure.",
            "Use this page as the opening summary before sharing the fuller supporting table.",
        ]
    return [
        "Confirm the business question this table is meant to answer.",
        "Add a sharper grouping before circulating the output broadly.",
        "Use the supporting table as backup, not as the lead story.",
    ]


def _chart_subtitle(chart_spec: dict, dimension_col: str, metric_col: str) -> str:
    chart_type = str(chart_spec.get("chart_type", "") or "").replace("_", " ").strip()
    if dimension_col and _looks_temporal(dimension_col) and metric_col:
        return f"{chart_type.title() if chart_type else 'Trend'} view of {_humanize(metric_col).lower()} over time."
    if chart_type and dimension_col and metric_col:
        return f"{chart_type.title()} view of {_humanize(metric_col).lower()} across {_humanize(dimension_col).lower()}."
    if chart_type:
        return f"{chart_type.title()} view selected for the clearest executive readout."
    return ""


def _humanize(value: str) -> str:
    label = re.sub(r"([a-z0-9])([A-Z])", r"\1 \2", str(value or ""))
    label = label.replace("_", " ")
    return " ".join(word.capitalize() for word in label.split())


def _format_value(column: str, value) -> str:
    try:
        numeric = float(value)
    except (TypeError, ValueError):
        return str(value)
    lowered = str(column or "").lower()
    if any(token in lowered for token in ("pct", "percent", "rate", "share")):
        return f"{numeric:,.1f}%"
    if any(token in lowered for token in ("salary", "income", "pay", "cost", "amount", "compensation")):
        return f"${numeric:,.0f}"
    if numeric.is_integer():
        return f"{numeric:,.0f}"
    return f"{numeric:,.1f}"


def _format_delta_value(column: str, value) -> str:
    try:
        numeric = float(value)
    except (TypeError, ValueError):
        return str(value)
    sign = "+" if numeric > 0 else "-" if numeric < 0 else ""
    lowered = str(column or "").lower()
    magnitude = abs(numeric)
    if any(token in lowered for token in ("pct", "percent", "rate", "share")):
        return f"{sign}{magnitude:,.1f} pts"
    if any(token in lowered for token in ("salary", "income", "pay", "cost", "amount", "compensation")):
        return f"{sign}${magnitude:,.0f}"
    if magnitude.is_integer():
        return f"{sign}{magnitude:,.0f}"
    return f"{sign}{magnitude:,.1f}"


def _format_period_point(value) -> str:
    if isinstance(value, pd.Timestamp):
        return value.strftime("%b %Y")
    return str(value)


def _looks_temporal(column: str) -> bool:
    lowered = str(column or "").lower()
    return any(token in lowered for token in ("date", "month", "quarter", "year", "week"))


def _time_series_frame(df: pd.DataFrame, dimension_col: str, metric_col: str) -> pd.DataFrame:
    if not dimension_col or not metric_col or dimension_col not in df.columns or metric_col not in df.columns:
        return pd.DataFrame()
    if not (_looks_temporal(dimension_col) or pd.api.types.is_datetime64_any_dtype(df[dimension_col])):
        return pd.DataFrame()

    frame = df.copy()
    frame[metric_col] = pd.to_numeric(frame[metric_col], errors="coerce")
    frame = frame.dropna(subset=[dimension_col, metric_col]).sort_values(dimension_col)
    return frame if not frame.empty else pd.DataFrame()


def render_chart_image(rows: list[dict], chart_spec: dict, *, width: int = 1180, height: int = 620) -> bytes | None:
    df = _prepare_dataframe(rows)
    if df.empty:
        return None

    chart_type = str(chart_spec.get("chart_type", "") or "").strip().lower()
    x_col = str(chart_spec.get("x_column", "") or "").strip()
    y_col = str(chart_spec.get("y_column", "") or "").strip()

    image = Image.new("RGB", (width, height), WHITE)
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((0, 0, width - 1, height - 1), radius=28, fill=WHITE, outline=LINE, width=2)

    if chart_type in {"lollipop", "bar", "horizontal_bar"} and x_col and y_col and x_col in df.columns and y_col in df.columns:
        _draw_rank_chart(draw, df, x_col, y_col, width, height, as_lollipop=(chart_type == "lollipop"))
    elif chart_type in {"line", "area"} and x_col and y_col and x_col in df.columns and y_col in df.columns:
        _draw_line_chart(draw, df, x_col, y_col, width, height, fill_area=(chart_type == "area"))
    elif chart_type in {"indicator", "summary"} and y_col and y_col in df.columns:
        _draw_indicator(draw, df, y_col, width, height)
    else:
        fallback_metric = _choose_metric_column(df)
        fallback_dimension = _choose_dimension_column(df, exclude={fallback_metric} if fallback_metric else None)
        if fallback_metric and fallback_dimension:
            _draw_rank_chart(draw, df, fallback_dimension, fallback_metric, width, height, as_lollipop=True)
        elif fallback_metric:
            _draw_indicator(draw, df, fallback_metric, width, height)
        else:
            _draw_empty_state(draw, width, height)

    buffer = BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def build_pdf_report(story: ReportStory) -> bytes:
    width, height = 1680, 960
    canvas = Image.new("RGB", (width, height), WHITE)
    draw = ImageDraw.Draw(canvas)

    draw.rectangle((0, 0, width, height), fill=WHITE)
    draw.rectangle((68, 72, 90, height - 72), fill=ACCENT)
    draw.rounded_rectangle((120, 72, width - 68, height - 72), radius=34, fill=(252, 253, 255), outline=LINE, width=2)

    draw.text((160, 104), story.title, font=_font(40, bold=True), fill=ACCENT_DARK)
    draw.text((160, 154), f"{story.scope_name} | {story.role} | {story.generated_label}", font=_font(17), fill=MUTED)

    headline_lines = _wrap_text(draw, story.headline, _font(56, bold=True), width - 290, max_lines=2)
    draw.multiline_text((160, 214), "\n".join(headline_lines), font=_font(56, bold=True), fill=INK, spacing=8)
    intro_y = 214 + (len(headline_lines) * 72) + 18
    sub_lines = _wrap_text(draw, story.subheadline, _font(24), width - 290, max_lines=3)
    draw.multiline_text((160, intro_y), "\n".join(sub_lines), font=_font(24), fill=MUTED, spacing=8)

    metric_y = intro_y + (len(sub_lines) * 34) + 42
    for index, metric in enumerate(story.key_metrics[:3]):
        left = 160 + (index * 432)
        _draw_metric_card(draw, metric, left, metric_y, 406, 148)

    chart_box = (160, metric_y + 190, 1038, height - 126)
    insight_box = (1076, metric_y + 190, width - 110, height - 126)
    draw.rounded_rectangle(chart_box, radius=28, fill=WHITE, outline=LINE, width=2)
    draw.rounded_rectangle(insight_box, radius=28, fill=PANEL, outline=LINE, width=2)

    draw.text((chart_box[0] + 24, chart_box[1] + 22), "PRIMARY VISUAL", font=_font(18, bold=True), fill=ACCENT)
    draw.text((chart_box[0] + 24, chart_box[1] + 54), story.chart_title, font=_font(28, bold=True), fill=INK)
    if story.chart_subtitle:
        chart_lines = _wrap_text(draw, story.chart_subtitle, _font(16), chart_box[2] - chart_box[0] - 48, max_lines=2)
        draw.multiline_text((chart_box[0] + 24, chart_box[1] + 92), "\n".join(chart_lines), font=_font(16), fill=MUTED, spacing=4)

    chart_top = chart_box[1] + 132
    if story.chart_image:
        chart_image = Image.open(BytesIO(story.chart_image)).convert("RGB")
        chart_image.thumbnail((chart_box[2] - chart_box[0] - 48, chart_box[3] - chart_top - 24))
        paste_x = chart_box[0] + ((chart_box[2] - chart_box[0]) - chart_image.width) // 2
        paste_y = chart_top + 6
        canvas.paste(chart_image, (paste_x, paste_y))
    else:
        _draw_empty_state(draw, chart_box[2] - chart_box[0] - 48, chart_box[3] - chart_top - 24, x=chart_box[0] + 24, y=chart_top)

    draw.text((insight_box[0] + 24, insight_box[1] + 22), "WHAT TO KNOW", font=_font(18, bold=True), fill=ACCENT)
    insight_y = insight_box[1] + 64
    for bullet in story.insights[:3]:
        insight_y = _draw_bullet(draw, bullet, x=insight_box[0] + 24, y=insight_y, width=insight_box[2] - insight_box[0] - 48, font=_font(22))

    draw.text((insight_box[0] + 24, insight_y + 8), "WHERE TO FOCUS", font=_font(18, bold=True), fill=ACCENT)
    action_y = insight_y + 52
    for bullet in story.actions[:3]:
        action_y = _draw_bullet(draw, bullet, x=insight_box[0] + 24, y=action_y, width=insight_box[2] - insight_box[0] - 48, font=_font(17), bullet_fill=SUCCESS, text_fill=MUTED)

    draw.line((160, height - 124, width - 110, height - 124), fill=LINE, width=2)
    footer_lines = _wrap_text(draw, story.source_note, _font(16), width - 300, max_lines=2)
    draw.multiline_text((160, height - 108), "\n".join(footer_lines), font=_font(16), fill=MUTED, spacing=4)

    output = BytesIO()
    canvas.save(output, format="PDF", resolution=180.0)
    return output.getvalue()


def build_ppt_report(story: ReportStory) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    slide = presentation.slides.add_slide(blank)
    _decorate_slide(slide)
    _slide_text(slide, 1.0, 0.55, 10.6, 0.3, story.title.upper(), 14, ACCENT, bold=True)
    _slide_text(slide, 1.0, 0.92, 11.0, 0.9, story.headline, 28, INK, bold=True)
    _slide_text(slide, 1.0, 1.82, 11.0, 0.5, story.subheadline, 14, MUTED)
    _slide_text(slide, 10.7, 0.58, 1.7, 0.2, story.generated_label, 11, MUTED, align=PP_ALIGN.RIGHT)
    for index, metric in enumerate(story.key_metrics[:3]):
        _slide_metric_card(slide, metric, 1.0 + (index * 3.95), 2.7)

    slide = presentation.slides.add_slide(blank)
    _decorate_slide(slide)
    _slide_text(slide, 1.0, 0.55, 6.0, 0.25, "CHART STORY", 14, ACCENT, bold=True)
    _slide_text(slide, 1.0, 0.9, 6.6, 0.45, story.chart_title, 24, INK, bold=True)
    _slide_text(slide, 1.0, 1.35, 6.6, 0.28, story.chart_subtitle, 12, MUTED)
    if story.chart_image:
        slide.shapes.add_picture(BytesIO(story.chart_image), Inches(0.95), Inches(1.95), width=Inches(7.35), height=Inches(4.65))
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.65), Inches(1.05), Inches(3.55), Inches(5.6))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(*PANEL)
    panel.line.color.rgb = RGBColor(*LINE)
    _slide_text(slide, 8.9, 1.3, 2.5, 0.2, "WHAT TO KNOW", 12, ACCENT, bold=True)
    y_cursor = 1.7
    for bullet in story.insights[:3]:
        y_cursor = _slide_bullet(slide, bullet, 8.9, y_cursor, 2.9, 0.58, 13, INK)
    _slide_text(slide, 8.9, y_cursor + 0.06, 2.5, 0.2, "WHERE TO FOCUS", 12, ACCENT, bold=True)
    y_cursor += 0.42
    for bullet in story.actions[:3]:
        y_cursor = _slide_bullet(slide, bullet, 8.9, y_cursor, 2.9, 0.48, 11, MUTED, bullet_color=SUCCESS)

    slide = presentation.slides.add_slide(blank)
    _decorate_slide(slide)
    _slide_text(slide, 1.0, 0.55, 8.5, 0.25, "SUPPORTING VIEW", 14, ACCENT, bold=True)
    _slide_text(slide, 1.0, 0.92, 8.5, 0.4, "Sample governed rows and export context", 20, INK, bold=True)
    _slide_text(slide, 1.0, 1.34, 8.5, 0.3, story.source_note, 11, MUTED)
    _slide_table(slide, story.preview_rows[:4], 1.0, 1.95, 7.6, 3.2)
    checklist = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.95), Inches(1.95), Inches(3.1), Inches(3.2))
    checklist.fill.solid()
    checklist.fill.fore_color.rgb = RGBColor(*ACCENT_SOFT)
    checklist.line.color.rgb = RGBColor(*LINE)
    _slide_text(slide, 9.2, 2.18, 2.2, 0.2, "ACTION CHECKLIST", 12, ACCENT, bold=True)
    y_cursor = 2.55
    for bullet in story.actions[:3]:
        y_cursor = _slide_bullet(slide, bullet, 9.18, y_cursor, 2.45, 0.52, 11, INK)

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _font(size: int, *, bold: bool = False):
    try:
        return ImageFont.truetype("DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf", size)
    except OSError:
        return ImageFont.load_default()


def _wrap_text(draw: ImageDraw.ImageDraw, text: str, font, max_width: int, *, max_lines: int = 3) -> list[str]:
    words = str(text or "").split()
    if not words:
        return [""]
    lines: list[str] = []
    current = words[0]
    for word in words[1:]:
        candidate = f"{current} {word}"
        if draw.textlength(candidate, font=font) <= max_width:
            current = candidate
        else:
            lines.append(current)
            current = word
    lines.append(current)
    if len(lines) > max_lines:
        lines = lines[:max_lines]
        lines[-1] = lines[-1].rstrip(".") + "..."
    return lines


def _draw_metric_card(draw, metric: ReportMetric, left: int, top: int, width: int, height: int) -> None:
    draw.rounded_rectangle((left, top, left + width, top + height), radius=24, fill=PANEL, outline=LINE, width=2)
    draw.text((left + 24, top + 22), metric.label.upper(), font=_font(15, bold=True), fill=ACCENT)
    draw.text((left + 24, top + 54), metric.value, font=_font(34, bold=True), fill=INK)
    if metric.note:
        lines = _wrap_text(draw, metric.note, _font(16), width - 48, max_lines=2)
        draw.multiline_text((left + 24, top + 102), "\n".join(lines), font=_font(16), fill=MUTED, spacing=4)


def _draw_bullet(draw, text: str, *, x: int, y: int, width: int, font, bullet_fill=ACCENT, text_fill=INK) -> int:
    draw.ellipse((x, y + 7, x + 12, y + 19), fill=bullet_fill)
    lines = _wrap_text(draw, text, font, width - 24, max_lines=4)
    draw.multiline_text((x + 22, y), "\n".join(lines), font=font, fill=text_fill, spacing=5)
    return y + (len(lines) * 30) + 14


def _draw_rank_chart(draw, df: pd.DataFrame, x_col: str, y_col: str, width: int, height: int, *, as_lollipop: bool) -> None:
    frame = df[[x_col, y_col]].dropna().copy().sort_values(y_col, ascending=False).head(8)
    if frame.empty:
        _draw_empty_state(draw, width, height)
        return
    left, top, right, bottom = 74, 72, width - 42, height - 62
    axis_left = left + 170
    axis_bottom = bottom - 22
    draw.line((axis_left, top + 10, axis_left, axis_bottom), fill=LINE, width=2)
    draw.line((axis_left, axis_bottom, right, axis_bottom), fill=LINE, width=2)
    max_value = float(pd.to_numeric(frame[y_col], errors="coerce").max()) or 1.0
    band_height = max(46, min(70, (axis_bottom - top - 18) // max(len(frame), 1)))
    for index, (_, row) in enumerate(frame.iterrows()):
        label = str(row[x_col])
        value = float(row[y_col])
        y = top + 14 + index * band_height
        draw.text((left, y + 6), label[:20], font=_font(16), fill=INK)
        usable_width = right - axis_left - 36
        bar_width = int((value / max_value) * usable_width)
        if as_lollipop:
            draw.line((axis_left + 12, y + 14, axis_left + 12 + bar_width, y + 14), fill=ACCENT, width=7)
            draw.ellipse((axis_left + 12 + bar_width - 11, y + 3, axis_left + 12 + bar_width + 11, y + 25), fill=ACCENT_DARK)
        else:
            draw.rounded_rectangle((axis_left + 12, y, axis_left + 12 + max(bar_width, 12), y + 28), radius=13, fill=ACCENT)
        value_text = _format_value(y_col, value)
        draw.text((min(right - 16 - draw.textlength(value_text, font=_font(16, bold=True)), axis_left + 20 + max(bar_width, 12)), y + 5), value_text, font=_font(16, bold=True), fill=INK)


def _draw_line_chart(draw, df: pd.DataFrame, x_col: str, y_col: str, width: int, height: int, *, fill_area: bool) -> None:
    frame = df[[x_col, y_col]].dropna().copy().sort_values(x_col).head(10)
    if frame.empty:
        _draw_empty_state(draw, width, height)
        return
    left, top, right, bottom = 82, 68, width - 46, height - 70
    draw.line((left, bottom, right, bottom), fill=LINE, width=2)
    draw.line((left, top, left, bottom), fill=LINE, width=2)
    values = pd.to_numeric(frame[y_col], errors="coerce")
    minimum = float(values.min())
    maximum = float(values.max())
    span = max(maximum - minimum, 1.0)
    step_x = (right - left - 34) / max(len(frame) - 1, 1)
    points = []
    for index, (_, row) in enumerate(frame.iterrows()):
        x = left + 18 + index * step_x
        y = bottom - 20 - ((float(row[y_col]) - minimum) / span) * (bottom - top - 54)
        points.append((x, y, row[x_col], row[y_col]))
    if fill_area:
        polygon = [(points[0][0], bottom - 1)] + [(point[0], point[1]) for point in points] + [(points[-1][0], bottom - 1)]
        draw.polygon(polygon, fill=(215, 232, 246))
    draw.line([(point[0], point[1]) for point in points], fill=ACCENT, width=6)
    for index, (x, y, label, value) in enumerate(points):
        draw.ellipse((x - 8, y - 8, x + 8, y + 8), fill=ACCENT_DARK)
        if index in {0, len(points) - 1} or len(points) <= 5:
            draw.text((x - 18, bottom + 10), str(label)[:10], font=_font(13), fill=MUTED)
        draw.text((x - 16, y - 26), _format_value(y_col, value), font=_font(13, bold=True), fill=INK)


def _draw_indicator(draw, df: pd.DataFrame, y_col: str, width: int, height: int) -> None:
    series = pd.to_numeric(df[y_col], errors="coerce").dropna()
    if series.empty:
        _draw_empty_state(draw, width, height)
        return
    value = series.iloc[0] if len(series) == 1 else series.mean()
    draw.rounded_rectangle((84, 82, width - 84, height - 82), radius=30, fill=PANEL, outline=LINE, width=2)
    draw.text((126, 126), _humanize(y_col).upper(), font=_font(18, bold=True), fill=ACCENT)
    draw.text((126, 216), _format_value(y_col, value), font=_font(74, bold=True), fill=INK)
    draw.text((126, 322), "Executive summary metric", font=_font(22), fill=MUTED)
    draw.line((126, 376, width - 126, 376), fill=LINE, width=2)
    draw.text((126, 418), "Use this view when the headline metric matters more than category comparison.", font=_font(18), fill=MUTED)


def _draw_empty_state(draw, width: int, height: int, *, x: int = 0, y: int = 0) -> None:
    draw.rounded_rectangle((x + 18, y + 18, x + width - 18, y + height - 18), radius=24, fill=PANEL, outline=LINE, width=2)
    draw.text((x + 46, y + (height // 2) - 14), "Chart preview unavailable for the current export.", font=_font(24, bold=True), fill=MUTED)


def _decorate_slide(slide) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.25), Inches(0.0), Inches(0.18), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(*ACCENT)
    band.line.fill.background()


def _slide_text(slide, left: float, top: float, width: float, height: float, text: str, size: int, color, *, bold: bool = False, align=PP_ALIGN.LEFT):
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = textbox.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text or "")
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    run.font.name = "Aptos"


def _slide_metric_card(slide, metric: ReportMetric, left: float, top: float) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.45), Inches(1.28))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*ACCENT_SOFT)
    box.line.color.rgb = RGBColor(*LINE)
    _slide_text(slide, left + 0.18, top + 0.12, 3.0, 0.18, metric.label.upper(), 10, ACCENT, bold=True)
    _slide_text(slide, left + 0.18, top + 0.42, 3.0, 0.36, metric.value, 23, INK, bold=True)
    _slide_text(slide, left + 0.18, top + 0.88, 3.0, 0.16, metric.note, 10, MUTED)


def _slide_bullet(slide, text: str, left: float, top: float, width: float, height: float, size: int, color, *, bullet_color=ACCENT) -> float:
    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left), Inches(top + 0.06), Inches(0.11), Inches(0.11))
    dot.fill.solid()
    dot.fill.fore_color.rgb = RGBColor(*bullet_color)
    dot.line.fill.background()
    _slide_text(slide, left + 0.18, top - 0.02, width, height, text, size, color)
    return top + height


def _slide_table(slide, rows: list[dict], left: float, top: float, width: float, height: float) -> None:
    if not rows:
        _slide_text(slide, left, top, width, 0.3, "No rows available for preview.", 12, MUTED)
        return
    columns = list(rows[0].keys())[:4]
    table = slide.shapes.add_table(len(rows) + 1, len(columns), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for column_index, column in enumerate(columns):
        cell = table.cell(0, column_index)
        cell.text = str(column)
        _style_slide_cell(cell, header=True)
    for row_index, row in enumerate(rows, start=1):
        for column_index, column in enumerate(columns):
            cell = table.cell(row_index, column_index)
            cell.text = str(row.get(column, ""))
            _style_slide_cell(cell, header=False)


def _style_slide_cell(cell, *, header: bool) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(*(ACCENT_SOFT if header else WHITE))
    paragraph = cell.text_frame.paragraphs[0]
    if paragraph.runs:
        run = paragraph.runs[0]
        run.font.size = Pt(10 if header else 9)
        run.font.bold = header
        run.font.color.rgb = RGBColor(*(ACCENT_DARK if header else INK))

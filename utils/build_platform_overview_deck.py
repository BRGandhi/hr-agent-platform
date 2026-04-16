from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT / "docs" / "HR_Insights_Platform_Overview_Deck.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(20, 46, 92)
BLUE = RGBColor(40, 79, 147)
INK = RGBColor(25, 36, 58)
SLATE = RGBColor(92, 110, 138)
PALE = RGBColor(244, 247, 252)
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(221, 179, 99)
MINT = RGBColor(196, 228, 216)
ROSE = RGBColor(241, 213, 213)
STEEL = RGBColor(214, 223, 236)


def set_slide_background(slide, color: RGBColor = PALE) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_soft_accent(slide, left: float, top: float, width: float, height: float, color: RGBColor, transparency: float = 0.12) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.fill.background()


def add_title_block(slide, title: str, subtitle: str = "", kicker: str = "") -> None:
    if kicker:
        kicker_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(3.0), Inches(0.35))
        tf = kicker_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = kicker.upper()
        p.font.name = "Aptos"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLUE

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.62), Inches(8.4), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = INK

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.28), Inches(8.9), Inches(0.7))
        tf = sub_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.name = "Aptos"
        p.font.size = Pt(15)
        p.font.color.rgb = SLATE


def style_text_frame(text_frame, font_name: str = "Aptos", font_size: int = 16, color: RGBColor = INK) -> None:
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color


def add_bullet_list(slide, items: list[str], left: float, top: float, width: float, height: float, font_size: int = 18) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    first = True
    for item in items:
        paragraph = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        paragraph.text = f"- {item}"
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = INK
        paragraph.space_after = Pt(6)


def add_card(slide, left: float, top: float, width: float, height: float, title: str, body: str, accent: RGBColor = BLUE) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = STEEL
    shape.line.width = Pt(1.1)

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.12), top + Inches(0.12), Inches(0.14), height - Inches(0.24))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(left + Inches(0.35), top + Inches(0.26), width - Inches(0.55), Inches(0.35))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Aptos"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = INK

    body_box = slide.shapes.add_textbox(left + Inches(0.35), top + Inches(0.68), width - Inches(0.55), height - Inches(0.88))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = body
    p.font.name = "Aptos"
    p.font.size = Pt(13)
    p.font.color.rgb = SLATE


def add_stat_card(slide, left: float, top: float, width: float, height: float, label: str, value: str, accent: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = STEEL
    shape.line.width = Pt(1.0)

    label_box = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.18), width - Inches(0.44), Inches(0.25))
    tf = label_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label.upper()
    p.font.name = "Aptos"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = SLATE

    value_box = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.46), width - Inches(0.44), Inches(0.5))
    tf = value_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = value
    p.font.name = "Aptos Display"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = accent


def add_footer(slide, page_num: int) -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(7.02), Inches(11.93), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = STEEL
    line.line.fill.background()

    box = slide.shapes.add_textbox(Inches(11.8), Inches(7.05), Inches(0.6), Inches(0.22))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = "Aptos"
    p.font.size = Pt(10)
    p.font.color.rgb = SLATE


def add_process_box(slide, left: float, top: float, width: float, title: str, body: str, accent: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(1.22))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.3)

    title_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.16), width - Inches(0.3), Inches(0.24))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Aptos"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = accent

    body_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.42), width - Inches(0.3), Inches(0.6))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = body
    p.font.name = "Aptos"
    p.font.size = Pt(11)
    p.font.color.rgb = SLATE


def add_arrow_text(slide, left: float, top: float) -> None:
    box = slide.shapes.add_textbox(left, top, Inches(0.35), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = ">"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos Display"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD


def build_deck() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_soft_accent(slide, Inches(8.9), Inches(0.35), Inches(3.4), Inches(1.4), GOLD, 0.16)
    add_soft_accent(slide, Inches(8.3), Inches(4.55), Inches(4.2), Inches(1.7), BLUE, 0.08)

    brand_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.68), Inches(0.55), Inches(2.2), Inches(0.42))
    brand_bar.fill.solid()
    brand_bar.fill.fore_color.rgb = NAVY
    brand_bar.line.fill.background()
    tf = brand_bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "HR Insights Platform"
    p.font.name = "Aptos"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.25), Inches(7.0), Inches(1.4))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "A governed, memory-aware HR analytics platform"
    p.font.name = "Aptos Display"
    p.font.size = Pt(29)
    p.font.bold = True
    p.font.color.rgb = INK

    sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(2.35), Inches(6.9), Inches(1.2))
    tf = sub_box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = (
        "Built for HR leaders and partners who need precise workforce answers, governed reports, "
        "polished visuals, and real continuity across prior work."
    )
    p.font.name = "Aptos"
    p.font.size = Pt(17)
    p.font.color.rgb = SLATE

    add_stat_card(slide, Inches(8.2), Inches(1.55), Inches(1.7), Inches(1.15), "Scope", "HR-only", BLUE)
    add_stat_card(slide, Inches(10.0), Inches(1.55), Inches(1.7), Inches(1.15), "Access", "Role-based", GOLD)
    add_stat_card(slide, Inches(8.2), Inches(2.9), Inches(1.7), Inches(1.15), "Memory", "Recall-first", BLUE)
    add_stat_card(slide, Inches(10.0), Inches(2.9), Inches(1.7), Inches(1.15), "UX", "Visual-ready", GOLD)

    add_bullet_list(
        slide,
        [
            "Answers HR-only questions with role-aware access controls",
            "Generates workforce reports and executive-ready visual options",
            "Preserves useful context across sessions, recalls saved insights, and explains methodology",
        ],
        Inches(0.78),
        Inches(4.05),
        Inches(7.2),
        Inches(1.7),
        font_size=17,
    )
    add_footer(slide, 1)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_title_block(slide, "Why this platform matters", "It is designed to feel like an HR analytics workspace, not a generic prompt box.", "Executive frame")
    add_card(slide, Inches(0.78), Inches(1.8), Inches(3.9), Inches(2.0), "Governed from the start", "Every request is checked against HR scope, approved business units, and allowed metric domains before meaningful model work begins.", BLUE)
    add_card(slide, Inches(4.88), Inches(1.8), Inches(3.9), Inches(2.0), "Built for decision support", "The platform answers questions, creates reports, recommends visuals, and helps users reuse prior insights without restarting their analysis.", GOLD)
    add_card(slide, Inches(8.98), Inches(1.8), Inches(3.55), Inches(2.0), "Enterprise-ready posture", "The architecture separates UI, policy, orchestration, tool execution, and data access so it can be hardened for regulated environments.", BLUE)

    add_card(slide, Inches(0.78), Inches(4.15), Inches(5.8), Inches(1.7), "Outcome for the user", "HR teams spend less time prompt-engineering and more time acting on workforce insights that are already grounded in access controls and prior context.", GOLD)
    add_card(slide, Inches(6.82), Inches(4.15), Inches(5.72), Inches(1.7), "Outcome for leadership", "The experience is explainable, reusable, and demo-ready: users can move from KPI question to report, chart, and saved recall in a single governed workflow.", BLUE)
    add_footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_title_block(slide, "Core feature stack", "The platform combines governed analytics, personalization, and polished output in one workflow.", "What it delivers")
    add_card(slide, Inches(0.78), Inches(1.8), Inches(2.85), Inches(1.8), "Governed HR Q&A", "Rejects non-HR asks, resolves role and coverage, and keeps answers inside approved workforce domains.", BLUE)
    add_card(slide, Inches(3.88), Inches(1.8), Inches(2.85), Inches(1.8), "Reports and exports", "Builds employee-level or aggregate reports and asks clarifying questions before generating underspecified outputs.", GOLD)
    add_card(slide, Inches(6.98), Inches(1.8), Inches(2.85), Inches(1.8), "Visual exploration", "Suggests multiple chart options, recommends the strongest one, and reuses the latest table context for visual follow-ups.", BLUE)
    add_card(slide, Inches(10.08), Inches(1.8), Inches(2.45), Inches(1.8), "Cached recall", "Reopens saved insights without rerunning the original SQL and lets the user continue naturally.", GOLD)
    add_card(slide, Inches(0.78), Inches(4.0), Inches(3.85), Inches(1.8), "Personalized workspace", "Headcount anchors the board, favorite topics stay visible, and the center panel adapts to prior HR activity.", GOLD)
    add_card(slide, Inches(4.88), Inches(4.0), Inches(3.85), Inches(1.8), "Helpful-answer feedback loop", "Yes/No feedback shapes favorite chats and reusable high-signal examples for future related asks.", BLUE)
    add_card(slide, Inches(8.98), Inches(4.0), Inches(3.55), Inches(1.8), "Explainable metrics", "Users can ask how a metric was calculated, which columns were used, and what snapshot caveats apply.", GOLD)
    add_footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_title_block(slide, "Why it feels cutting edge", "The platform is not one prompt template. It behaves like a governed analytic agent.", "Differentiators")
    add_bullet_list(
        slide,
        [
            "Agentic routing before tool use: each turn is classified into modes such as data_query, report, policy, history_lookup, or visual_follow_up.",
            "Multi-layer context assembly: access profile, recent memory, helpful prior answers, context documents, and latest-table state are merged before the answer loop.",
            "Strict relevance instead of noisy memory stuffing: relevant-chat suggestions only surface for strong close matches.",
            "Recall-first UX: saved insights reopen as cached summaries and immediately seed the active session for natural follow-up work.",
            "Methodology explanations are in-scope: users can ask how a result was calculated and receive a grounded explanation of definition, columns, formula, and caveats.",
        ],
        Inches(0.82),
        Inches(1.8),
        Inches(7.25),
        Inches(4.6),
        font_size=17,
    )
    add_card(slide, Inches(8.55), Inches(1.95), Inches(3.95), Inches(1.55), "Not just a chatbot", "The agent chooses tools, reuses memory, validates access, and keeps the live analytic thread intact.", BLUE)
    add_card(slide, Inches(8.55), Inches(3.78), Inches(3.95), Inches(1.55), "Not just BI", "The user can ask in natural language, get clarifications when needed, and continue with saved or generated results.", GOLD)
    add_footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_title_block(slide, "Governed by design", "Access control is not bolted on later. It shapes the workflow from the first turn.", "Trust model")
    add_process_box(slide, Inches(0.8), Inches(2.15), Inches(2.05), "User + role", "Signed-in identity resolves role, business coverage, allowed metrics, and document tags.", BLUE)
    add_arrow_text(slide, Inches(2.95), Inches(2.52))
    add_process_box(slide, Inches(3.25), Inches(2.15), Inches(2.05), "Orchestrator", "Checks HR scope, resolves follow-up context, and routes the request before model work begins.", GOLD)
    add_arrow_text(slide, Inches(5.4), Inches(2.52))
    add_process_box(slide, Inches(5.7), Inches(2.15), Inches(2.05), "Governed tools", "Queries, reports, visuals, and retrieval tools run only inside approved HR boundaries.", BLUE)
    add_arrow_text(slide, Inches(7.85), Inches(2.52))
    add_process_box(slide, Inches(8.15), Inches(2.15), Inches(2.05), "Context layer", "Policy docs, metric definitions, memory, and latest-table context are injected narrowly and intentionally.", GOLD)
    add_arrow_text(slide, Inches(10.3), Inches(2.52))
    add_process_box(slide, Inches(10.6), Inches(2.15), Inches(1.95), "Answer", "Returns a role-aware response, with follow-up prompts and saved-memory support.", BLUE)

    add_card(slide, Inches(1.1), Inches(4.45), Inches(5.2), Inches(1.45), "What this prevents", "Out-of-scope questions, unauthorized metrics, weakly related history suggestions, and accidental data leakage beyond approved business units.", GOLD)
    add_card(slide, Inches(6.6), Inches(4.45), Inches(5.25), Inches(1.45), "What this enables", "A trustworthy HR experience where governance is visible, useful, and still lightweight enough for day-to-day analytics work.", BLUE)
    add_footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_title_block(slide, "Personalization and navigation", "The home screen and sidebar behave like a memory-aware HR workspace.", "User experience")
    add_card(slide, Inches(0.8), Inches(1.85), Inches(3.8), Inches(1.8), "Center board", "Headcount shows first when available. Other KPI families appear only after the user has actually explored them. Remaining slots become concrete next-question cards.", BLUE)
    add_card(slide, Inches(4.78), Inches(1.85), Inches(3.8), Inches(1.8), "Favorite Topics", "The HR themes a user revisits most often stay visible and open by default, giving the workspace a personalized starting point.", GOLD)
    add_card(slide, Inches(8.76), Inches(1.85), Inches(3.75), Inches(1.8), "Favorite / Relevant / Past Chats", "Favorite chats are shaped by reuse and feedback, relevant chats are intentionally strict, and past chats preserve the broader history.", BLUE)
    add_card(slide, Inches(0.8), Inches(4.1), Inches(5.85), Inches(1.65), "Substantive questions stay featured", "Thin shorthand follow-ups like yes, show me, or answer question 1 are filtered out of featured-history surfaces so the board keeps showing the real business question.", GOLD)
    add_card(slide, Inches(6.88), Inches(4.1), Inches(5.62), Inches(1.65), "Immediate refresh loop", "Feedback actions and new chats refresh the personalized history state so the workspace evolves inside the same browser session.", BLUE)
    add_footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_title_block(slide, "Follow-up intelligence and explainability", "The platform maintains thread continuity instead of forcing users to restate their work.", "Conversation design")
    add_card(slide, Inches(0.8), Inches(1.9), Inches(3.85), Inches(1.75), "Short follow-ups", "Replies like yes, break it down, job level, or show me inherit the active HR thread and latest meaningful anchor.", BLUE)
    add_card(slide, Inches(4.85), Inches(1.9), Inches(3.85), Inches(1.75), "Clarification before output", "If a report or table request is underspecified, the assistant asks for missing columns or cuts before generation.", GOLD)
    add_card(slide, Inches(8.9), Inches(1.9), Inches(3.6), Inches(1.75), "Metric explanation requests", "Users can ask how a result was calculated, what columns were used, and what the formula means without tripping the guardrail.", BLUE)
    add_bullet_list(
        slide,
        [
            "Example: 'How many employees were promoted in the last year?' followed by 'show me how you calculated this metric' should return the definition, columns, formula, and snapshot caveat.",
            "If prior context is missing, the platform asks which HR metric or prior result to explain instead of refusing the request.",
            "This makes the agent more useful in real HR conversations, where trust often depends on showing the work behind the answer.",
        ],
        Inches(0.92),
        Inches(4.18),
        Inches(11.2),
        Inches(1.75),
        font_size=16,
    )
    add_footer(slide, 7)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_title_block(slide, "Reporting and visualization workflow", "The platform moves smoothly from question to table to chart to explanation.", "Analytic flow")
    add_process_box(slide, Inches(0.85), Inches(2.0), Inches(2.2), "Ask", "Natural-language HR question, report request, or follow-up.", BLUE)
    add_arrow_text(slide, Inches(3.1), Inches(2.37))
    add_process_box(slide, Inches(3.35), Inches(2.0), Inches(2.25), "Retrieve / query", "The agent calls the right SQL or retrieval tool with scope-aware constraints.", GOLD)
    add_arrow_text(slide, Inches(5.68), Inches(2.37))
    add_process_box(slide, Inches(5.95), Inches(2.0), Inches(2.2), "Render", "Tables, markdown sections, charts, and saved-memory cards appear in one consistent surface.", BLUE)
    add_arrow_text(slide, Inches(8.2), Inches(2.37))
    add_process_box(slide, Inches(8.45), Inches(2.0), Inches(2.25), "Recommend", "Visual options compare chart shapes and recommend the strongest story quickly.", GOLD)
    add_arrow_text(slide, Inches(10.78), Inches(2.37))
    add_process_box(slide, Inches(11.02), Inches(2.0), Inches(1.45), "Reuse", "Latest-table context stays live for charting and explanation follow-ups.", BLUE)

    add_card(slide, Inches(1.15), Inches(4.35), Inches(5.2), Inches(1.55), "Why this matters in a demo", "You can show a clean path from KPI question to departmental breakdown, chart recommendation, and methodology explanation without resetting the thread.", GOLD)
    add_card(slide, Inches(6.75), Inches(4.35), Inches(5.05), Inches(1.55), "Why this matters in production", "The same flow supports repeated HR analysis over time because outputs, feedback, and recalled insights keep feeding the personalization layer.", BLUE)
    add_footer(slide, 8)

    # Slide 9
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_title_block(slide, "Architecture at a glance", "A clean separation of UI, orchestration, tools, memory, and data access keeps the system extensible and governable.", "System view")
    add_process_box(slide, Inches(0.72), Inches(1.95), Inches(1.7), "Browser UI", "Sidebar, KPI board, chat surface, chart actions, feedback loop", BLUE)
    add_arrow_text(slide, Inches(2.5), Inches(2.33))
    add_process_box(slide, Inches(2.8), Inches(1.95), Inches(1.9), "FastAPI", "Auth-aware endpoints for stats, history, chat, recall, feedback, export", GOLD)
    add_arrow_text(slide, Inches(4.82), Inches(2.33))
    add_process_box(slide, Inches(5.12), Inches(1.95), Inches(2.0), "Orchestrator", "Routing, context assembly, access checks, tool loop, final answer", BLUE)
    add_arrow_text(slide, Inches(7.23), Inches(2.33))
    add_process_box(slide, Inches(7.52), Inches(1.95), Inches(1.9), "Tool layer", "SQL, reports, memory retrieval, context docs, visual suggestion", GOLD)
    add_arrow_text(slide, Inches(9.52), Inches(2.33))
    add_process_box(slide, Inches(9.82), Inches(1.95), Inches(2.65), "Data + models", "hr_data.db, access_control.db, context_store.db, Anthropic or OpenAI-compatible LLMs", BLUE)

    add_card(slide, Inches(1.05), Inches(4.38), Inches(3.45), Inches(1.5), "Memory plane", "context_store.db supports recent memory, strong-match retrieval, feedback-weighted favorites, and cached recall.", BLUE)
    add_card(slide, Inches(4.88), Inches(4.38), Inches(3.45), Inches(1.5), "Policy plane", "Access profiles and document tags constrain what questions, metrics, and documents are in scope for a user.", GOLD)
    add_card(slide, Inches(8.72), Inches(4.38), Inches(3.55), Inches(1.5), "Experience plane", "The UI makes the agent feel coherent: history, charts, reports, and explanations all connect as one workflow.", BLUE)
    add_footer(slide, 9)

    # Slide 10
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide)
    add_title_block(slide, "Suggested demo story", "A short flow that shows the platform's core strengths in one pass.", "Walkthrough")
    add_bullet_list(
        slide,
        [
            "Start on the personalized workspace: call out headcount first, favorite topics, and the role / coverage / access chips.",
            "Ask a direct HR question: 'What is the total headcount for my business units?'",
            "Use a short follow-up: 'break it down by department' or 'job role' to show continuity.",
            "Turn the result into a chart: 'build me a chart of department and attrition rate' or use Visual options.",
            "Ask an explanation follow-up: 'show me how you calculated this metric' to highlight explainability.",
            "Open a saved chat from the sidebar and continue from recalled context to show recall without rerun.",
        ],
        Inches(0.85),
        Inches(1.85),
        Inches(7.1),
        Inches(4.8),
        font_size=18,
    )
    add_card(slide, Inches(8.45), Inches(2.0), Inches(4.0), Inches(1.35), "Best prompt categories", "Headcount, attrition, promotions, satisfaction, access-envelope questions, and chart follow-ups.", BLUE)
    add_card(slide, Inches(8.45), Inches(3.65), Inches(4.0), Inches(1.35), "Avoid in demo", "Month-over-month claims, rolling 12-month trend language, or real employee-name expectations in the IBM sample data.", GOLD)
    add_card(slide, Inches(8.45), Inches(5.3), Inches(4.0), Inches(1.35), "Closing message", "This is a governed HR analytics assistant with memory, explainability, and a polished workflow - not a generic chatbot.", BLUE)
    add_footer(slide, 10)

    # Slide 11
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, WHITE)
    add_soft_accent(slide, Inches(0.85), Inches(1.25), Inches(5.7), Inches(4.7), BLUE, 0.07)
    add_soft_accent(slide, Inches(7.2), Inches(1.65), Inches(4.8), Inches(3.9), GOLD, 0.11)
    add_title_block(slide, "Why this platform stands out", "A modern HR analytics platform should be governed, contextual, explainable, and actually pleasant to use.", "Closing")
    add_card(slide, Inches(0.95), Inches(2.0), Inches(3.7), Inches(1.55), "Governed", "Role-based access, approved business coverage, metric-domain limits, and protected retrieval keep the experience safe by default.", BLUE)
    add_card(slide, Inches(0.95), Inches(3.8), Inches(3.7), Inches(1.55), "Contextual", "Short replies, recalled chats, and methodology questions all continue from the right HR anchor instead of resetting the conversation.", GOLD)
    add_card(slide, Inches(4.95), Inches(2.0), Inches(3.7), Inches(1.55), "Explainable", "Users can inspect definitions, columns, formulas, and snapshot caveats behind the outputs they are asked to trust.", GOLD)
    add_card(slide, Inches(4.95), Inches(3.8), Inches(3.7), Inches(1.55), "Reusable", "Feedback, favorites, relevant chats, and past chats turn every good answer into future product value.", BLUE)
    add_card(slide, Inches(8.95), Inches(2.55), Inches(3.2), Inches(2.1), "Cutting edge in practice", "It combines agentic routing, retrieval, memory, visual reuse, and explainability in one governed HR workflow that is ready to demo and extend.", BLUE)
    add_footer(slide, 11)

    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    build_deck()
    print(f"Wrote {OUTPUT_PATH}")
